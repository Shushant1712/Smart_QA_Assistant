from pptx import Presentation


def extract_pptx(file_path: str):
    presentation = Presentation(file_path)
    extracted_data = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        if slide_text:
            extracted_data.append({
                "text": "\n".join(slide_text),
                "metadata": {
                    "slide": slide_num,
                    "source": file_path
                }
            })

    return extracted_data
